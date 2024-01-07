
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
import PIL.Image as Image
from pathlib import Path
import pandas as pd
import warnings
import math
import io
import re
import numpy as np

QMD_Head = '''---
format:
  revealjs:
    theme: THEME
---

{{< include _imports.qmd >}}
'''

Slide = '''
TITLE
:::: {.columns}

::: {.column width="WIDTH1%"}

TEXT

:::

::: {.column width="WIDTH2%"}

DATA

:::

::::
'''

Img = '''
![](FN)
'''

Tbl = '''
```{r}
#| label: tbl-LABEL
#| tbl-cap: CAPTION

Data <- read.csv(file = 'FN',sep=',')

kable(Data, col.names = gsub("[.]", " ", names(Data)))

```
'''

Fig = '''
```{r}
#| label: fig-LABEL
#| fig-cap: CAPTION

Data <- read.csv(file = 'FN',sep=',', check.names = FALSE)

fig <- plot_ly(data=Data)
TRACES
fig

```
'''

Traces = {
'Scatter':'''
fig <- fig %>% add_trace(x = ~X, y = ~~y, name = '~y',type = "scatter", mode = "markers")
''',
'Scatter_Lines':'''
fig <- fig %>% add_trace(x = ~X, y = ~~y, name = '~y',type = "scatter", mode = "lines+markers")
''',
'Line':'''
fig <- fig %>% add_trace(x = ~X, y = ~~y, name = '~y',type = "scatter",mode = "lines")
''',
'Bar':'''
fig <- fig %>% add_trace(x = ~X, y = ~~y,  name = '~y', type = "bar")
''',
}

Notes = '''

::: {.notes}

NOTES

:::

'''

class QMD():
    def __init__(self,in_file,out_dir,out_name,theme='default',trouble_shoot = False):
        self.QMD = QMD_Head.replace('THEME',theme)
        self.report = {
            'slide_number':[],
            'issue':[],
            'Title':[]
        }
        prs = Presentation(in_file)
        P_file = Parse(out_dir)
        for i,slide in enumerate(prs.slides):
            if trouble_shoot == True:
                print(i+1)
            P_file.ParseSlide(slide,i+1,out_name)
            self.QMD += P_file.FullSlide
            self.report['slide_number'].append(i+1)
            self.report['issue'].append(P_file.issues)
            self.report['Title'].append(P_file.Title)
        
        with open(f'{out_dir}{out_name}.qmd','w', encoding="utf-8") as out:
            out.write(self.QMD)
        
class UnGroup():
    def __init__(self,shapes):
        self.shapes = [shape for shape in shapes]
        self.AllShapes = []
        while len(self.shapes)>0:
            self.UnGroup()

    def UnGroup(self):
        shape = self.shapes[0]
        self.CheckGroup(shape)

    def CheckGroup(self,shape):
        if shape.shape_type == 6:
            self.shapes.pop(0)
            for s in shape.shapes:
                self.shapes.append(s)
        else:
            self.AllShapes.append(shape)
            self.shapes.pop(0)

class Parse():
    
    def __init__(self,dir=''):
        self.dir=dir
        self.Slide = Slide
        self.Fig = Fig
        self.Tbl = Tbl
        self.Img = Img
        self.Notes = Notes
        self.Traces = Traces
        self.maxdim=1750
        self.Title = ''
        
        self.codes=pd.read_csv('MSO_Type_Codes.csv')
        self.chart_codes=pd.read_csv('MSO_Chart_Codes.csv')

    def ParseSlide(self,slide,n,LN):
        self.n = n
        self.LN = LN
        self.Clear()
        Shapes = UnGroup(slide.shapes).AllShapes
        for i,shape in enumerate(Shapes):
            self.Handle_Shape(shape,f'{n}_{i}')
        self.Text_Frames = pd.DataFrame(self.text_frames)
        if self.Text_Frames.shape[0]==0:
            self.WriteTitle('')
        else:
            for i,rank in enumerate(np.sort(self.Text_Frames['fontSize'].unique())[::-1]):
                self.Text_Frames.loc[self.Text_Frames['fontSize']==rank,'Order'] = i
            self.Text_Frames['Order'] = self.Text_Frames['Order'].fillna(i)
            for i,rank in enumerate(np.sort(self.Text_Frames['len'].unique())):
                self.Text_Frames.loc[self.Text_Frames['len']==rank,'Order'] += i
            for i,rank in enumerate(np.sort(self.Text_Frames['top'].unique())[::-1]):
                self.Text_Frames.loc[self.Text_Frames['top']==rank,'Order'] += i
            for i,rank in enumerate(np.sort(self.Text_Frames['left'].unique())):
                self.Text_Frames.loc[self.Text_Frames['left']==rank,'Order'] += i/2
            self.Text_Frames.loc[self.Text_Frames['text'].str.len()<=1,'Order']*=10
            self.Text_Frames = self.Text_Frames.set_index('Order').sort_index()
            for i,text in enumerate(self.Text_Frames['text']):
                if i == 0:
                    self.WriteTitle(text)
                else:
                   self.WriteLine(text)
        self.FullSlide = self.FullSlide.replace('TEXT',self.Text)
        self.FullSlide = self.FullSlide.replace('DATA',self.Data_Text)
        
        if len(self.Text)>0 and len(self.Data_Text)>0:
            self.FullSlide = self.FullSlide.replace('WIDTH1','50').replace('WIDTH2','50')
        elif len(self.Text)>0 and len(self.Data_Text)==0:
            self.FullSlide = self.FullSlide.replace('WIDTH1','100').replace('WIDTH2','0')
        else:
            self.FullSlide = self.FullSlide.replace('WIDTH1','0').replace('WIDTH2','100')
            
        if slide.has_notes_slide == True:
            self.FullSlide += self.Notes.replace('NOTES',slide.notes_slide.notes_text_frame.text)

                
    def Clear(self):
        self.Text = ''
        self.Data_Text = ''
        self.text_frames = {
            'text':[],
            'type':[],
            'top':[],
            'left':[],
            'fontSize':[],
            'len':[]}
        self.issues = ''


    def Handle_Shape(self,shape,name): 
        try:
            if (shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX) or (shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.has_text_frame):       
                self.GetText(shape)
            elif (shape.shape_type == MSO_SHAPE_TYPE.PICTURE) or (shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, "image")):
                issue = self.WriteImage(shape,name)
                if len(issue) > 0:
                    self.issues += f"{issue}|"
            elif (shape.shape_type == MSO_SHAPE_TYPE.TABLE) or (shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.has_table):
                self.MakeTable(shape,name)
            elif (shape.shape_type == MSO_SHAPE_TYPE.CHART) or (shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.has_chart):
                self.MakeChart(shape,name)
            else:
                type = self.codes.loc[self.codes['Value']==shape.shape_type,'Description'].values[0]
                self.issues += f"{type}:{shape.name}|"
        except:
            self.issues += f"Other Error:{shape.name}|"
            pass

    def GetText(self,shape):
        self.text_frames['fontSize'].append(shape.text_frame.paragraphs[0].font.size)
        self.text_frames['text'].append(shape.text)
        self.text_frames['type'].append(shape.shape_type)
        self.text_frames['top'].append(shape.top)
        self.text_frames['left'].append(shape.left)
        self.text_frames['len'].append(len(shape.text))
        
    def WriteTitle(self,str):
        str = str.split('|')[-1].lstrip()
        if self.n == 1:
            self.Title = str
        if len(str)<=2:
            self.FullSlide = self.Slide.replace('TITLE','\n\n---\n\n')
        else:
            self.FullSlide = self.Slide.replace('TITLE',f'\n\n## {str}\n\n')

    def WriteLine(self,str):
        spl = str.split('\n')
        if len(spl)>1:
            pref='* '
            suff='\n'
        elif len(spl)==1:
            pref='\n'
            suff='\n\n'
        for sub in spl:
            if len(sub)>2:
                self.Text += f'{pref}{sub}{suff}'

    def WriteImage(self,shape,name):
        shape_Name = shape.name
        try:
            self.type = shape.image.ext.lower()
            self.image = Image.open(io.BytesIO(shape.image.blob))
            root = f"{self.dir}images/{self.LN}"
            Path(root).mkdir(parents=True, exist_ok=True)
            if self.type!='jpg' and self.type!='png' and self.type!='gif':
                self.image = self.image.convert('RGB')
                self.fn = f"{root}/{name}.jpg"
            else:
                self.fn = f"{root}/{name}.{self.type}"
            if self.image.size[0] > self.maxdim and self.image.size[1]<self.image.size[0]:
                wpercent = (self.maxdim/float(self.image.size[0]))
                hsize = int((float(self.image.size[1])*float(wpercent)))
                self.image = self.image.resize((self.maxdim,hsize), resample=Image.BICUBIC)
            elif self.image.size[1] > self.maxdim and self.image.size[1]>self.image.size[0]:
                wpercent = (self.maxdim/float(self.image.size[1]))
                wsize = int((float(self.image.size[0])*float(wpercent)))
                self.image = self.image.resize((wsize,self.maxdim), resample=Image.BICUBIC)
            self.image.save(self.fn,quality=95,optimize=True)
            self.Data_Text += '\n\n' +self.Img.replace('FN',self.relDir(self.fn)).replace('LABEL',name).replace('CAPTION','')
            return('')
        except:
            return(f'Failed to Render:{shape_Name}')
        
    def MakeTable(self,shape,name):
        root = f"{self.dir}Data/{self.LN}"
        Path(root).mkdir(parents=True, exist_ok=True)
        self.fn = f"{root}/{name}_Table.csv"
        data = shape.table
        table = {}
        for i,row in enumerate(data.rows):
            for j, cell in enumerate(row.cells):
                if i == 0:
                    table[j] = [cell.text]
                else:
                    table[j].append(cell.text)

        self.Table = pd.DataFrame(data=table)
        self.Table.to_csv(self.fn,index=False)
        self.Data_Text += '\n\n' +self.Tbl.replace('FN',self.relDir(self.fn)).replace('LABEL',name)

    def MakeChart(self,shape,name=''):
        
        self.chart = shape.chart
        self.part = shape.chart_part
        self.type = self.chart_codes.loc[self.chart_codes['Value']==self.chart.chart_type,'Type'].values[0]

        root = f"{self.dir}Data/{self.LN}"
        Path(root).mkdir(parents=True, exist_ok=True)
        self.fn = f"{root}/{name}_{self.type}.csv"
        self.formatTable()
        self.Data_Text += '\n\n' +self.Fig_Traces.replace('FN',self.relDir(self.fn)).replace('LABEL',name)

    def formatTable(self):
        
        blob = self.part.chart_workbook.xlsx_part.blob
        workbook_xml = io.BytesIO(blob)
        workbook_xml.seek(0)
        with warnings.catch_warnings(record=True):
            warnings.simplefilter("always")
            self.Table = pd.read_excel(workbook_xml,engine='openpyxl')
        Series = [s.name for s in self.chart.series]
        self.Names = [re.sub(r'\W+', '', s.name) for s in self.chart.series]
        if len([i for i in Series if i in self.Table.columns.values]) == 0:
            self.Table = self.Table.set_index(self.Table.columns[0]).T
            self.Table = self.Table.reset_index()
        self.Table = self.Table.dropna(how='all', axis=1)
        self.Table = self.Table.dropna(how='all', axis=0)
                
        Drop = []
        Thresh = 5
        for col in self.Table.columns:
            if self.Table.iloc[:Thresh][col].isna().sum()>Thresh-1:
                Drop.append(col)
        self.Table.drop(columns=Drop)
        

        RN = {}
        for col in self.Table.columns:
            try:
                if self.type != 'Bar':
                    self.Table[col] = self.Table[col].str.extract(r"([+-]?(?:\d*\.*\d+))", expand=False)
            except:
                pass
            RN[col] = col
        for name in Series:
            RN[name] = re.sub(r'\W+', '', RN[name])+'.y'
        b = self.Names[0]
        for col in self.Table.columns:
            if 'Unnamed' in col or 'index' in col:
                RN[col] = b+'.x'
            else:
                b = re.sub(r'\W+', '', col)
        self.Table = self.Table.rename(columns=RN)


        for name in self.Names:
            if name+'.x' not in self.Table.columns:
                X = [x for x in self.Table.columns if 'x' in x]
                self.Table[name+'.x'] = self.Table[X[-1]].copy()

        Dups = self.Table.loc[:,self.Table.apply(lambda x: x.duplicated(keep = False),axis=1).all()].columns
        if len(Dups)>1:
            self.Table = self.Table.drop(Dups[1:],axis=1)
            self.Table = self.Table.rename(columns={Dups[0]:'X'})
            newY = {}
            for col in self.Table.columns:
                s = col.split('.')
                if len(s)>1:
                    newY[col] = s[0]
                
            self.Table = self.Table.rename(columns=newY)
            self.addTraces()
        else:
            self.addTraces(mode = 'VariableX')
        self.Table.to_csv(self.fn,index=False)
        
    def addTraces(self,mode='OneX'):
        Traces = ''
        for trace in self.Names:
            if mode == 'OneX':
                Traces += self.Traces[self.type].replace('~y',trace)
            elif mode == 'VariableX':
                Traces += self.Traces[self.type].replace('~y',f'{trace}.y').replace('~X',f'~{trace}.x')

        self.Fig_Traces = self.Fig.replace('TRACES',Traces)

    def relDir(self,pth):
        return(pth.replace(self.dir,''))
    